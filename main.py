from fastapi import FastAPI, UploadFile, File
from fastapi.responses import JSONResponse
import fitz
import os
import tempfile
import uuid
import requests
import base64
from paddleocr import PaddleOCR
import signal
import sys
import cloudinary
import cloudinary.uploader
import cloudinary.api
from dotenv import load_dotenv
import concurrent.futures

# Load environment variables
load_dotenv()

# Windows Ctrl+C handling
def signal_handler(sig, frame):
    print('You pressed Ctrl+C! Exiting Python server...')
    sys.exit(0)

signal.signal(signal.SIGINT, signal_handler)

app = FastAPI()

# Configure Cloudinary
cloudinary.config(
  cloud_name = os.getenv("CLOUDINARY_CLOUD_NAME"),
  api_key = os.getenv("CLOUDINARY_API_KEY"),
  api_secret = os.getenv("CLOUDINARY_API_SECRET"),
  secure = True
)

ocr_instance = None

def get_ocr():
    global ocr_instance
    if ocr_instance is None:
        print("Loading PaddleOCR model (Bulletproof CPU Mode)...")
        try:
            # Disable MKLDNN and force CPU for maximum compatibility on Windows
            # show_log=False reduces terminal noise and potential output stream issues
            ocr_instance = PaddleOCR(
                use_angle_cls=False, 
                lang="en", 
                use_gpu=False, 
                enable_mkldnn=False, 
                show_log=False, 
                ocr_version='PP-OCRv4'
            )
            print("PaddleOCR loaded successfully.")
        except Exception as e:
            print(f"CRITICAL ERROR: PaddleOCR failed to initialize: {e}")
            ocr_instance = None
    return ocr_instance

def google_vision_ocr(file_path):
    """Performs OCR using Google Cloud Vision API REST endpoint."""
    api_key = os.getenv("GOOGLE_SEARCH_API_KEY")
    if not api_key:
        print("ERROR: GOOGLE_SEARCH_API_KEY not found for Vision OCR.")
        return ""
        
    try:
        print(f"DEBUG: Running Google Vision OCR on {file_path}")
        
        # 1. Image Check
        if not os.path.exists(file_path):
            print(f"CRITICAL: file_path {file_path} does not exist!")
            return ""
            
        file_size = os.path.getsize(file_path)
        print(f"DEBUG: Image file exists. Size: {file_size} bytes")
        if file_size == 0:
            print("CRITICAL: Image file is empty (0 bytes)! MuPDF rendering likely failed.")
            return ""

        with open(file_path, "rb") as image_file:
            content = base64.b64encode(image_file.read()).decode("utf-8")

        url = f"https://vision.googleapis.com/v1/images:annotate?key={api_key}"
        payload = {
            "requests": [
                {
                    "image": {"content": content},
                    "features": [{"type": "DOCUMENT_TEXT_DETECTION"}]
                }
            ]
        }
        
        print(f"DEBUG: Sending POST to Vision API... (Project: {os.getenv('GOOGLE_SEARCH_API_KEY', '')[:10]}...)")
        response = requests.post(url, json=payload, timeout=60)
        
        print(f"DEBUG: Vision API Header response: {response.status_code}")
        if response.status_code != 200:
            print(f"ERROR: Google Vision API returned HTTP {response.status_code}")
            print(f"Response Body: {response.text}")
            return ""

        data = response.json()
        responses = data.get("responses", [])
        if not responses:
            print("ERROR: Empty 'responses' array from Google Vision.")
            return ""
            
        res0 = responses[0]
        if "error" in res0:
            err_msg = res0['error'].get('message', 'Unknown error')
            print(f"ERROR: Google Vision Result Error: {err_msg}")
            # Log the full error for debugging
            import json
            print(f"Full Error Detail: {json.dumps(res0['error'])}")
            return ""
            
        text = res0.get("fullTextAnnotation", {}).get("text", "")
        print(f"DEBUG: Google Vision OCR result success! Length: {len(text)}")
        return text
    except Exception as e:
        print(f"Google Vision OCR exception: {e}")
        import traceback
        traceback.print_exc()
        return ""

def perform_ocr_on_file(file_path):
    """Helper to run OCR on a file path directly. 
    Uses Google Vision as primary, PaddleOCR as fallback."""
    try:
        # 1. Try Google Vision (High Accuracy)
        text = google_vision_ocr(file_path)
        if text and len(text.strip()) > 5:
            return text
            
        # 2. Fallback to PaddleOCR (Local)
        print("Falling back to PaddleOCR...")
        ocr = get_ocr()
        if ocr is None:
            return ""
            
        print(f"DEBUG: Running PaddleOCR on {file_path}")
        result = ocr.ocr(file_path)
        
        text = ""
        if result and result[0]:
            for line in result[0]:
                if line and len(line) > 1:
                    text += line[1][0] + "\n"
        
        print(f"DEBUG: OCR result length for {file_path}: {len(text)}")
        return text
    except Exception as e:
        print(f"OCR failed for {file_path}: {str(e)}")
        # If it's the MKLDNN error specifically, we log a warning for the user
        if "MKLDNN" in str(e) or "ConvertPirAttribute" in str(e):
             print("⚠️  Warning: PaddleOCR Windows compatibility issue detected. Falling back to next engine...")
        return f"[Python OCR failed for this page: {str(e)}]"

def upload_to_cloudinary(file_path):
    """Uploads a file to Cloudinary and returns the secure URL."""
    try:
        if not os.getenv("CLOUDINARY_CLOUD_NAME"):
            print("WARNING: Cloudinary not configured. Skipping upload.")
            return None
            
        print(f"Uploading {file_path} to Cloudinary...")
        response = cloudinary.uploader.upload(file_path, folder="velocity_ai_uploads")
        url = response.get("secure_url")
        print(f"Uploaded to Cloudinary: {url}")
        return url
    except Exception as e:
        print(f"Cloudinary upload failed: {e}")
        return None

def process_pptx(content):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import io
    
    print(f"Processing PPTX content...")
    ppt_file = io.BytesIO(content)
    prs = Presentation(ppt_file)
    text = ""
    
    for i, slide in enumerate(prs.slides):
        text += f"\n--- Slide {i+1} ---\n"
        # 0. Speaker Notes (Important for context)
        if slide.has_notes_slide:
            try:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes = notes_slide.notes_text_frame.text
                    if notes.strip():
                        text += f"\n[Speaker Notes]: {notes.strip()}\n"
            except:
                pass

        for shape in slide.shapes:
            # 1. Text extraction
            if hasattr(shape, "text"):
                text += shape.text + "\n"

            # 2. Table extraction
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                try:
                    table = shape.table
                    # Build Markdown Table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text_frame.text.strip().replace('\n', ' ') for cell in row.cells]
                        rows.append(f"| {' | '.join(cells)} |")
                    
                    if rows:
                        # Header separator
                        separator = f"| {' | '.join(['---'] * len(table.columns))} |"
                        rows.insert(1, separator)
                        text += "\n" + "\n".join(rows) + "\n"
                except Exception as e:
                    print(f"Error extracting table: {e}")
            
            # 3. Chart extraction
            if hasattr(shape, "has_chart") and shape.has_chart:
                try:
                    chart = shape.chart
                    title = "Untitled Chart"
                    try:
                        if chart.chart_title and chart.chart_title.has_text_frame:
                            title = chart.chart_title.text_frame.text
                    except:
                        pass
                        
                    text += f"\n**Chart Data: {title}**\n"
                    
                    # Try to build a table of the chart data
                    for plot in chart.plots:
                        try:
                            categories = [c.label for c in plot.categories]
                            series_list = plot.series
                            
                            # Markdown table for chart data
                            # Headers: Category | Series 1 | Series 2 ...
                            headers = ["Category"] + [s.name for s in series_list]
                            rows = [f"| {' | '.join(headers)} |"]
                            rows.append(f"| {' | '.join(['---'] * len(headers))} |")
                            
                            for i, category in enumerate(categories):
                                row_data = [str(category)]
                                for series in series_list:
                                    try:
                                        val = series.values[i]
                                        row_data.append(str(val) if val is not None else "N/A")
                                    except:
                                        row_data.append("N/A")
                                rows.append(f"| {' | '.join(row_data)} |")
                            
                            text += "\n".join(rows) + "\n"
                        except:
                            continue
                except Exception as e:
                    print(f"Error extracting chart data: {e}")

            # 4. Image extraction (Upload to Cloudinary + OCR)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_blob = shape.image.blob
                    
                    # Create temp file for the image
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
                        tmp_img.write(image_blob)
                        image_path = tmp_img.name
                        
                    try:
                        # Upload to Cloudinary
                        image_url = upload_to_cloudinary(image_path)
                        if image_url:
                            # Add Markdown Image Link using Cloudinary URL
                            text += f"\n![Image]({image_url})\n"
                        else:
                            text += f"\n[Image upload failed]\n"
                        
                        # Also run OCR for accessibility/context
                        try:
                            # CENTRALIZED: Using perform_ocr_on_file instead of direct PaddleOCR
                            img_text = perform_ocr_on_file(image_path)
                            
                            if img_text.strip():
                                text += f"\n[Image Text/OCR]: {img_text.strip()}\n"
                        except Exception as ocr_err:
                            print(f"OCR failed for image: {ocr_err}")
                            
                    finally:
                        # Clean up temp file
                        if os.path.exists(image_path):
                            os.remove(image_path)
                            
                except Exception as e:
                    print(f"Error extracting image: {e}")
    return text

@app.post("/extract")
def extract(file: UploadFile = File(...)):
    temp_pdf_path = None
    try:
        content = file.file.read() 
        name = (file.filename or "").lower()
        
        if name.endswith(".pdf"):
            doc = None
            try:
                # Try opening from memory first (Fastest)
                doc = fitz.open(stream=content, filetype="pdf")
                total_pages = len(doc)
                print(f"DEBUG: PDF opened from memory. Total pages: {total_pages}")
            except Exception as e:
                # Disk fallback for Windows stability (Fixes Errno 22)
                print(f"DEBUG: Memory load failed, using disk fallback...")
                with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
                    tmp_pdf.write(content)
                    temp_pdf_path = tmp_pdf.name
                doc = fitz.open(temp_pdf_path)
                total_pages = len(doc)

            if total_pages == 0:
                if doc: doc.close()
                return { "text": "[Empty PDF]" }

            # --- PARALLEL PROCESSING HELPER ---
            def process_single_page(page_num):
                # We open a NEW handle for each thread to avoid 'fitz' multithreading issues
                handle = None
                try:
                    if temp_pdf_path:
                        handle = fitz.open(temp_pdf_path)
                    else:
                        handle = fitz.open(stream=content, filetype="pdf")
                    
                    p = handle[page_num]
                    
                    # 1. Direct text extraction (Fastest)
                    t = p.get_text().strip()
                    if len(t) > 200:
                        handle.close()
                        return f"\n--- Page {page_num+1} ---\n{t}"
                    
                    # 2. OCR Fallback (150 DPI for optimal Speed + Accuracy)
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                        img_path = tmp.name
                    
                    pix = p.get_pixmap(dpi=150)
                    pix.save(img_path)
                    handle.close()
                    handle = None
                    
                    print(f"DEBUG: Parallel OCR'ing Page {page_num+1}...")
                    page_text = perform_ocr_on_file(img_path)
                    
                    if os.path.exists(img_path):
                        os.remove(img_path)
                        
                    return f"\n--- Page {page_num+1} ---\n{page_text}"
                except Exception as e:
                    if handle:
                        try: handle.close()
                        except: pass
                    return f"\n--- Page {page_num+1} ERROR: {str(e)} ---"

            # Parallel execution (5 workers: stable on Windows, avoids network congestion)
            print(f"DEBUG: Starting parallel extraction for {total_pages} pages...")
            with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
                results = list(executor.map(process_single_page, range(total_pages)))
            
            doc.close()
            return { "text": "\n".join(results) }

        elif name.endswith(".pptx"):
            print(f"DEBUG: Processing PPTX file: {name}")
            return { "text": process_pptx(content) }

        else:
            # Image file processing
            print(f"DEBUG: Processing Image file: {name}")
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(content)
                img_path = tmp.name
            
            text = perform_ocr_on_file(img_path)
            if os.path.exists(img_path): os.remove(img_path)
            return { "text": text }

    except Exception as e:
        print(f"Extraction Error for {file.filename}: {str(e)}")
        import traceback
        traceback.print_exc()
        return JSONResponse(status_code=500, content={"error": str(e)})
    finally:
        if temp_pdf_path and os.path.exists(temp_pdf_path):
            try: os.remove(temp_pdf_path)
            except: pass

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000)